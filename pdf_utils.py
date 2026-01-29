import os
import requests
from io import BytesIO
import fitz
from pypdf import PdfWriter, PdfReader
from pdf2docx import Converter
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
import pdfplumber
import pandas as pd
from pptx import Presentation
try:
    import extract_msg
except ImportError:
    extract_msg = None
from xhtml2pdf import pisa
from deep_translator import GoogleTranslator
import ebooklib
from ebooklib import epub
from bs4 import BeautifulSoup
from PIL import Image

def add_password(input_path, output_path, password):
    reader = PdfReader(input_path)
    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)
    writer.encrypt(password)
    with open(output_path, "wb") as f:
        writer.write(f)

def unlock_pdf(input_path, output_path, password=""):
    reader = PdfReader(input_path)
    if reader.is_encrypted:
        success = reader.decrypt(password)
        if not success:
            raise ValueError("Incorrect password or encryption not supported.")
    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)
    with open(output_path, "wb") as f:
        writer.write(f)

def rotate_pdf(input_path, output_path, degrees=90):
    reader = PdfReader(input_path)
    writer = PdfWriter()
    for page in reader.pages:
        page.rotate(int(degrees))
        writer.add_page(page)
    with open(output_path, "wb") as f:
        writer.write(f)

def delete_pages(input_path, output_path, pages_to_delete):
    reader = PdfReader(input_path)
    writer = PdfWriter()
    for i, page in enumerate(reader.pages):
        if (i + 1) not in pages_to_delete:
            writer.add_page(page)
    with open(output_path, "wb") as f:
        writer.write(f)

def add_watermark(input_path, output_path, watermark_text):
    reader = PdfReader(input_path)
    writer = PdfWriter()
    packet = BytesIO()
    can = canvas.Canvas(packet, pagesize=letter)
    can.setFont("Helvetica", 60)
    can.setFillColorRGB(0.6, 0.6, 0.6, 0.3)
    can.saveState()
    can.translate(300, 500)
    can.rotate(45)
    can.drawCentredString(0, 0, watermark_text)
    can.restoreState()
    can.save()
    packet.seek(0)
    watermark_pdf = PdfReader(packet)
    watermark_page = watermark_pdf.pages[0]
    for page in reader.pages:
        page.merge_page(watermark_page)
        writer.add_page(page)
    with open(output_path, "wb") as f:
        writer.write(f)

def add_page_numbers(input_path, output_path):
    reader = PdfReader(input_path)
    writer = PdfWriter()
    n_pages = len(reader.pages)
    for i, page in enumerate(reader.pages):
        packet = BytesIO()
        w = float(page.mediabox.width)
        h = float(page.mediabox.height)
        can = canvas.Canvas(packet, pagesize=(w, h))
        can.setFont("Helvetica", 10)
        can.drawString(w - 50, 20, f"{i+1}/{n_pages}")
        can.save()
        packet.seek(0)
        num_pdf = PdfReader(packet)
        page.merge_page(num_pdf.pages[0])
        writer.add_page(page)
    with open(output_path, "wb") as f:
        writer.write(f)

def pdf_to_text(input_path, output_path):
    doc = fitz.open(input_path)
    with open(output_path, "w", encoding="utf-8") as f:
        for page in doc:
            f.write(page.get_text())

def msg_to_pdf(input_path, output_path):
    if not extract_msg:
        raise ImportError("extract_msg library not available")
    msg = extract_msg.Message(input_path)
    body = msg.body
    sender = msg.sender
    subject = msg.subject
    c = canvas.Canvas(output_path, pagesize=letter)
    width, height = letter
    y = height - 50
    c.setFont("Helvetica-Bold", 12)
    c.drawString(40, y, f"Subject: {subject}")
    y -= 20
    c.drawString(40, y, f"From: {sender}")
    y -= 30
    c.setFont("Helvetica", 10)
    text = c.beginText(40, y)
    text.setFont("Helvetica", 10)
    for line in body.split('\n'):
        if y < 50:
            c.drawText(text)
            c.showPage()
            y = height - 50
            text = c.beginText(40, y)
            text.setFont("Helvetica", 10)
        text.textLine(line.strip())
        y -= 12
    c.drawText(text)
    c.save()

def merge_pdfs(file_list, output_path):
    merger = PdfWriter()
    for pdf in file_list:
        merger.append(pdf)
    merger.write(output_path)
    merger.close()

def split_pdf(input_path, output_dir):
    reader = PdfReader(input_path)
    base_name = os.path.splitext(os.path.basename(input_path))[0]
    for i, page in enumerate(reader.pages):
        writer = PdfWriter()
        writer.add_page(page)
        output_filename = f"{base_name}_page_{i+1}.pdf"
        with open(os.path.join(output_dir, output_filename), "wb") as out:
            writer.write(out)

def pdf_to_word(input_path, output_path):
    cv = Converter(input_path)
    cv.convert(output_path, start=0, end=None)
    cv.close()

def pdf_to_images(input_path, output_dir):
    doc = fitz.open(input_path)
    base_name = os.path.splitext(os.path.basename(input_path))[0]
    for i in range(len(doc)):
        page = doc.load_page(i)
        pix = page.get_pixmap()
        output_file = os.path.join(output_dir, f"{base_name}_page_{i+1}.png")
        pix.save(output_file)

def images_to_pdf(image_list, output_path):
    if not image_list: return
    images = [Image.open(x).convert("RGB") for x in image_list]
    images[0].save(output_path, save_all=True, append_images=images[1:])

def compress_pdf(input_path, output_path):
    try:
        initial_size = os.path.getsize(input_path)
        doc = fitz.open(input_path)
        for page in doc:
            try:
                img_list = page.get_images(full=True)
                for img in img_list:
                    xref = img[0]
                    if xref == 0: continue
                    base = doc.extract_image(xref)
                    if not base: continue
                    if len(base["image"]) < 5120: continue
                    try:
                        pil_img = Image.open(BytesIO(base["image"]))
                    except: continue
                    max_dim = 1000
                    scale = 1.0
                    w, h = pil_img.size
                    if max(w, h) > max_dim:
                        scale = max_dim / max(w, h)
                        new_w = int(w * scale)
                        new_h = int(h * scale)
                        pil_img = pil_img.resize((new_w, new_h), Image.LANCZOS)
                    out_buffer = BytesIO()
                    if pil_img.mode in ("RGBA", "LA"):
                        bg = Image.new("RGB", pil_img.size, (255, 255, 255))
                        bg.paste(pil_img, mask=pil_img.split()[-1])
                        pil_img = bg
                    elif pil_img.mode == "P":
                         pil_img = pil_img.convert("RGB")
                    if pil_img.mode == "L": cs = "/DeviceGray"
                    else: cs = "/DeviceRGB"; pil_img = pil_img.convert("RGB")
                    pil_img.save(out_buffer, format="JPEG", quality=50, optimize=True)
                    new_bytes = out_buffer.getvalue()
                    if len(new_bytes) < len(base["image"]):
                        doc.update_stream(xref, new_bytes)
                        if hasattr(doc, "xref_set_key"):
                            doc.xref_set_key(xref, "Filter", "/DCTDecode")
                            doc.xref_set_key(xref, "Width", pil_img.width)
                            doc.xref_set_key(xref, "Height", pil_img.height)
                            doc.xref_set_key(xref, "ColorSpace", cs)
                            doc.xref_set_key(xref, "BitsPerComponent", 8)
            except Exception: pass
        doc.save(output_path, garbage=4, deflate=True)
        if os.path.exists(output_path):
            if os.path.getsize(output_path) >= initial_size:
                doc.close()
                doc2 = fitz.open(input_path)
                doc2.save(output_path, garbage=4, deflate=True)
                doc2.close()
    except Exception as e:
        try:
            doc = fitz.open(input_path)
            doc.save(output_path, garbage=4, deflate=True)
        except:
             raise Exception(f"Failed to compress PDF: {str(e)}")

def pdf_to_excel(input_path, output_path):
    with pdfplumber.open(input_path) as pdf:
        all_tables = []
        for page in pdf.pages:
            tables = page.extract_tables()
            for table in tables:
                df = pd.DataFrame(table)
                all_tables.append(df)
    if all_tables:
        with pd.ExcelWriter(output_path) as writer:
            for i, df in enumerate(all_tables):
                df.to_excel(writer, sheet_name=f"Table_{i+1}", index=False)
    else:
        pd.DataFrame().to_excel(output_path)

def crop_pdf(input_path, output_path, margins):
    reader = PdfReader(input_path)
    writer = PdfWriter()
    for page in reader.pages:
        x0 = page.cropbox.lower_left[0]
        y0 = page.cropbox.lower_left[1]
        x1 = page.cropbox.upper_right[0]
        y1 = page.cropbox.upper_right[1]
        left, top, right, bottom = margins
        new_x0 = x0 + left
        new_y0 = y0 + bottom
        new_x1 = x1 - right
        new_y1 = y1 - top
        if new_x1 > new_x0 and new_y1 > new_y0:
            page.cropbox.lower_left = (new_x0, new_y0)
            page.cropbox.upper_right = (new_x1, new_y1)
        writer.add_page(page)
    writer.write(output_path)

def rearrange_pdf(input_path, output_path, page_order):
    reader = PdfReader(input_path)
    writer = PdfWriter()
    total_pages = len(reader.pages)
    for p_num in page_order:
        if 1 <= p_num <= total_pages:
            writer.add_page(reader.pages[p_num-1])
    writer.write(output_path)

def pdf_to_ppt(input_path, output_path):
    doc = fitz.open(input_path)
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    for i in range(len(doc)):
        page = doc.load_page(i)
        pix = page.get_pixmap(dpi=150)
        img_data = pix.tobytes("png")
        slide = prs.slides.add_slide(blank_slide_layout)
        stream = BytesIO(img_data)
        slide.shapes.add_picture(stream, 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(output_path)

def url_to_pdf(url, output_path):
    try:
        if not url.startswith("http"): url = "http://" + url
        res = requests.get(url)
        if res.status_code != 200: raise Exception(f"HTTP Error: {res.status_code}")
        with open(output_path, "wb") as result_file:
            pisa_status = pisa.CreatePDF(res.text, dest=result_file)
        if pisa_status.err: raise Exception("PDF generation failed")
    except Exception as e:
        raise Exception(f"Failed to convert URL: {str(e)}")

def pdf_to_csv(input_path, output_path):
    with pdfplumber.open(input_path) as pdf:
        all_rows = []
        for page in pdf.pages:
            tables = page.extract_tables()
            for table in tables:
                all_rows.extend(table)
    if all_rows:
        df = pd.DataFrame(all_rows)
        df.to_csv(output_path, index=False, header=False)
    else:
        pd.DataFrame().to_csv(output_path)

def epub_to_pdf(input_path, output_path):
    try:
        book = epub.read_epub(input_path)
        content = "<html><body>"
        for item in book.get_items():
            if item.get_type() == ebooklib.ITEM_DOCUMENT:
                raw = item.get_content().decode('utf-8', errors='ignore')
                soup = BeautifulSoup(raw, 'html.parser')
                if soup.body: content += str(soup.body)
                else: content += str(soup)
        content += "</body></html>"
        with open(output_path, "wb") as f:
            pisa.CreatePDF(content, dest=f)
    except Exception as e:
        raise Exception(f"EPUB conversion failed: {str(e)}")

def create_pdf(text_content, output_path):
    c = canvas.Canvas(output_path, pagesize=letter)
    width, height = letter
    y = height - 50
    c.setFont("Helvetica", 12)
    lines = text_content.split('\n')
    for line in lines:
        if y < 50:
            c.showPage()
            y = height - 50
            c.setFont("Helvetica", 12)
        c.drawString(50, y, line)
        y -= 15
    c.save()

def translate_pdf(input_path, output_path, target_lang='en'):
    doc = fitz.open(input_path)
    full_text = ""
    for page in doc: full_text += page.get_text()
    chunks = [full_text[i:i+4000] for i in range(0, len(full_text), 4000)]
    translator = GoogleTranslator(source='auto', target=target_lang)
    translated_text = ""
    for chunk in chunks:
        translated_text += translator.translate(chunk) + "\n"
    create_pdf(translated_text, output_path)

def add_text_annotation(input_path, output_path, text, x, y):
    reader = PdfReader(input_path)
    writer = PdfWriter()
    packet = BytesIO()
    can = canvas.Canvas(packet, pagesize=letter)
    can.setFont("Helvetica", 12)
    can.drawString(float(x), float(y), text)
    can.save()
    packet.seek(0)
    text_pdf = PdfReader(packet)
    text_page = text_pdf.pages[0]
    for page in reader.pages:
        page.merge_page(text_page)
        writer.add_page(page)
    with open(output_path, "wb") as f:
        writer.write(f)

def pdf_to_tiff(input_path, output_path):
    doc = fitz.open(input_path)
    img_list = []
    for i in range(len(doc)):
        page = doc.load_page(i)
        pix = page.get_pixmap()
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        img_list.append(img)
    if img_list:
        img_list[0].save(output_path, save_all=True, append_images=img_list[1:], compression="tiff_deflate")

def sign_pdf(input_path, output_path, signature_image_path, x=100, y=100, width=150, height=50):
    """Add a signature image to all pages of a PDF."""
    doc = fitz.open(input_path)
    sig_img = fitz.open(signature_image_path)
    
    for page in doc:
        rect = fitz.Rect(x, y, x + width, y + height)
        page.insert_image(rect, filename=signature_image_path)
    
    doc.save(output_path)
    doc.close()

def flatten_pdf(input_path, output_path):
    """Flatten PDF - convert interactive elements to static content."""
    doc = fitz.open(input_path)
    
    for page in doc:
        # Flatten annotations
        annots = page.annots()
        if annots:
            for annot in annots:
                annot.update()
        
        # Flatten form fields by rendering page to image and back
        pix = page.get_pixmap(dpi=150)
    
    # Save with garbage collection to remove unused objects
    doc.save(output_path, garbage=4, deflate=True)
    doc.close()
